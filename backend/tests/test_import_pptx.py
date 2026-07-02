"""문서 가져오기 — PPTX → 보고서 초안(휴리스틱). python-pptx 로 인메모리 .pptx 를
만들어 업로드 → 제목/문단/표/이미지 위젯이 extra_blocks 로 생성되는지 검증.
"""
from __future__ import annotations

import base64
from io import BytesIO

from fastapi.testclient import TestClient

from app.database import SessionLocal
from app.main import app
from app.modules.auth.services import create_access_token

ADMIN = 2

# 1x1 PNG.
_PNG = base64.b64decode(
    "iVBORw0KGgoAAAANSUhEUgAAAAEAAAABCAQAAAC1HAwCAAAAC0lEQVR42mNkYPhfDwAChwGA60e6kgAAAABJRU5ErkJggg=="
)


def _h(uid, slug):
    return {
        "Authorization": f"Bearer {create_access_token(uid)}",
        "X-Workspace-Slug": slug,
    }


def _make_pptx() -> bytes:
    from pptx import Presentation
    from pptx.util import Inches

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # Title Only
    slide.shapes.title.text = "제목 슬라이드"
    tb = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(5), Inches(1))
    tf = tb.text_frame
    tf.text = "첫 문단"
    tf.add_paragraph().text = "둘째 문단"
    table = slide.shapes.add_table(2, 2, Inches(1), Inches(4), Inches(5), Inches(1)).table
    table.cell(0, 0).text = "이름"
    table.cell(0, 1).text = "값"
    table.cell(1, 0).text = "철수"
    table.cell(1, 1).text = "90"
    slide.shapes.add_picture(BytesIO(_PNG), Inches(1), Inches(6), Inches(1), Inches(1))
    buf = BytesIO()
    prs.save(buf)
    return buf.getvalue()


def _drop_files(file_ids):
    if not file_ids:
        return
    from app.modules.files.models import File

    db = SessionLocal()
    try:
        for fid in file_ids:
            f = db.get(File, fid)
            if f:
                db.delete(f)
        db.commit()
    finally:
        db.close()


def test_import_pptx_returns_draft():
    """PPTX → report_archive_draft_v1 draft(보고서 미생성). 4종 위젯·이미지 file_id 확인."""
    c = TestClient(app)
    data = _make_pptx()
    r = c.post(
        "/api/imports/pptx",
        headers=_h(ADMIN, "personal-2"),
        files={
            "file": (
                "테스트.pptx",
                data,
                "application/vnd.openxmlformats-officedocument.presentationml.presentation",
            )
        },
    )
    assert r.status_code in (200, 201), r.text
    draft = r.json()["data"]["draft"]
    assert draft["_type"] == "report_archive_draft_v1"
    page = draft["pages"][0]
    extra = page.get("extra_blocks") or []
    content = page.get("content") or {}
    types = {b.get("type") for b in extra}
    assert {"heading", "rich_text", "table", "image"} <= types, types
    img_ids = [b["id"] for b in extra if b.get("type") == "image"]
    assert img_ids
    files = (content.get(img_ids[0]) or {}).get("files") or []
    assert files and files[0].get("file_id"), content.get(img_ids[0])
    assert page.get("name") == "제목 슬라이드"
    _drop_files([f["file_id"] for f in files])


def test_parse_pptx_overlay_fills_merged_form_table():
    """표지 양식 패턴: 표는 라벨+빈 병합 값칸, 값은 위에 겹친 텍스트박스.
    좌표 매칭으로 값칸을 채우고 텍스트박스는 별도 문단으로 남기지 않아야 한다.
    """
    from pptx import Presentation
    from pptx.util import Inches, Pt

    from app.modules.imports.pptx_parser import parse_pptx

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
    # 1행 3열 표: [라벨][값(병합 2칸, 빈칸)]
    tbl = slide.shapes.add_table(
        1, 3, Inches(0), Inches(0), Inches(6), Inches(0.5)
    ).table
    for i in range(3):
        tbl.columns[i].width = Inches(2)  # c0:0-2, c1:2-4, c2:4-6
    tbl.cell(0, 0).text = "과제명"
    tbl.cell(0, 1).merge(tbl.cell(0, 2))  # 값칸 병합(빈칸)
    # 값 텍스트박스를 병합 값칸(c1) 위에 겹쳐 놓는다.
    b1 = slide.shapes.add_textbox(Inches(2.1), Inches(0.05), Inches(3), Inches(0.4))
    b1.text_frame.text = "낙하 시뮬레이션"
    b1.text_frame.paragraphs[0].font.size = Pt(12)

    buf = BytesIO()
    prs.save(buf)
    res = parse_pptx(buf.getvalue())
    seg_types = [s["type"] for s in res["pages"][0]["segments"]]
    assert seg_types == ["table"], seg_types  # 텍스트박스가 문단으로 새지 않음
    t = res["pages"][0]["segments"][0]["content"]
    flat = [c["label"] for c in t["columns"]] + [
        v for row in t["rows"] for v in row.values()
    ]
    assert "낙하 시뮬레이션" in flat, flat
    assert "과제명" in flat, flat
    assert "열 2" not in flat and "열 3" not in flat, flat  # 병합 빈칸 → 채워짐


def test_parse_pptx_reads_embedded_excel_table():
    """엑셀에서 붙여넣은 임베드(OLE, .xlsx) 표를 표 세그먼트로 읽는다."""
    import openpyxl
    from pptx import Presentation
    from pptx.enum.shapes import PROG_ID
    from pptx.util import Inches

    from app.modules.imports.pptx_parser import parse_pptx

    wb = openpyxl.Workbook()
    ws = wb.active
    ws.append(["이름", "점수"])
    ws.append(["철수", 90])
    ws.append(["영희", 85])
    xbuf = BytesIO()
    wb.save(xbuf)
    xbuf.seek(0)

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.shapes.add_ole_object(
        xbuf,
        PROG_ID.XLSX,
        left=Inches(1),
        top=Inches(1),
        width=Inches(4),
        height=Inches(2),
    )
    pbuf = BytesIO()
    prs.save(pbuf)

    res = parse_pptx(pbuf.getvalue())
    segs = res["pages"][0]["segments"]
    tables = [s for s in segs if s["type"] == "table"]
    assert tables, [s["type"] for s in segs]
    t = tables[0]["content"]
    assert [c["label"] for c in t["columns"]] == ["이름", "점수"], t["columns"]
    body = [tuple(r[c["key"]] for c in t["columns"]) for r in t["rows"]]
    assert ("철수", "90") in body, body
    assert ("영희", "85") in body, body


def test_import_pptx_requires_auth():
    c = TestClient(app)
    r = c.post(
        "/api/imports/pptx",
        files={"file": ("x.pptx", b"not-a-pptx", "application/octet-stream")},
    )
    assert r.status_code in (401, 403), r.text
