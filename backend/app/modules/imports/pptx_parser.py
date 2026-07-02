"""PPTX → 위젯 세그먼트(휴리스틱, DB 무관·순수).

슬라이드 1개 = 보고서 페이지 1개. 슬라이드 shape 를 (top,left) 순으로 훑어
제목→heading, 텍스트프레임→rich_text, 표→table, 그림→image 세그먼트로 만든다.
차트·SmartArt 등 미지원은 건너뛰고 warnings 에 남긴다(무손실 가장 금지).

반환: {"pages": [{"name", "segments":[...]}], "warnings": [...]}
  세그먼트: {type:'heading', content:{text,level}}
          | {type:'rich_text', content:{items:[{depth,text}]}}
          | {type:'table', content:{columns:[{key,label,type}], rows:[{col_N:val}]}}
          | {type:'image', image_bytes, image_mime, image_ext, alt}
이미지 업로드(→file_id)와 draft 조립은 라우트가 한다(여기선 바이트만 담아 반환).
"""
from __future__ import annotations

from io import BytesIO

MAX_DEPTH = 5


def _pos(shape):
    """(top,left) EMU. None 은 맨 위/왼쪽으로."""
    top = getattr(shape, "top", None)
    left = getattr(shape, "left", None)
    return (top if top is not None else -1, left if left is not None else -1)


def _flatten(shapes):
    """그룹 shape 를 펼쳐 평평한 리스트로."""
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    out = []
    for s in shapes:
        is_group = False
        try:
            is_group = s.shape_type == MSO_SHAPE_TYPE.GROUP
        except Exception:
            is_group = False
        if is_group:
            try:
                out.extend(_flatten(s.shapes))
                continue
            except Exception:
                pass
        out.append(s)
    return out


def _table_segment(rows):
    rows = [r for r in rows if any((c or "").strip() for c in r)] or rows
    if not rows:
        return None
    ncol = max(len(r) for r in rows)
    header = rows[0]
    columns = [
        {
            "key": f"col_{i + 1}",
            "label": (header[i].strip() if i < len(header) and header[i] else "") or f"열 {i + 1}",
            "type": "text",
        }
        for i in range(ncol)
    ]
    body = []
    for r in rows[1:]:
        obj = {}
        for i, col in enumerate(columns):
            obj[col["key"]] = (r[i].strip() if i < len(r) and r[i] else "")
        body.append(obj)
    return {"type": "table", "content": {"columns": columns, "rows": body}}


def _table_grid(tbl, overlay=None):
    """표를 문자열 격자로. 병합 원점(origin)의 텍스트를 span 영역 전체로 채운다
    (python-pptx 는 원점에만 텍스트를 주므로 그대로 두면 양식형 표가 빈칸투성이가 됨).
    overlay={(r,c): text} 는 표 위에 겹쳐 놓은 텍스트박스에서 뽑은 값칸 내용으로,
    비어 있는 셀에만 채운다(병합 값칸이면 span 전체로).
    """
    nr, nc = len(tbl.rows), len(tbl.columns)
    grid = [["" for _ in range(nc)] for _ in range(nr)]

    def put(r, c, txt):
        if not txt:
            return
        cell = tbl.cell(r, c)
        try:
            origin = cell.is_merge_origin
        except Exception:
            origin = False
        if origin:
            try:
                sh, sw = cell.span_height, cell.span_width
            except Exception:
                sh = sw = 1
            for dr in range(sh or 1):
                for dc in range(sw or 1):
                    if r + dr < nr and c + dc < nc and not grid[r + dr][c + dc]:
                        grid[r + dr][c + dc] = txt
        elif not grid[r][c]:
            grid[r][c] = txt

    # 원본 셀 텍스트(라벨)
    for r in range(nr):
        for c in range(nc):
            put(r, c, (tbl.cell(r, c).text or "").strip())
    # 겹쳐 놓은 값
    for (r, c), txt in (overlay or {}).items():
        if 0 <= r < nr and 0 <= c < nc:
            put(r, c, txt)
    return grid


def _table_geom(shape):
    """표 셀 격자의 누적 좌표(행 top 경계 · 열 left 경계, EMU)."""
    tbl = shape.table
    top = shape.top or 0
    left = shape.left or 0
    row_tops = [top]
    for row in tbl.rows:
        row_tops.append(row_tops[-1] + (row.height or 0))
    col_lefts = [left]
    for col in tbl.columns:
        col_lefts.append(col_lefts[-1] + (col.width or 0))
    return row_tops, col_lefts


def _locate_cell(row_tops, col_lefts, top, left):
    """(top,left) 점이 들어가는 (행,열). 격자 밖이면 None."""
    r = c = None
    for i in range(len(row_tops) - 1):
        if row_tops[i] <= top < row_tops[i + 1]:
            r = i
            break
    for j in range(len(col_lefts) - 1):
        if col_lefts[j] <= left < col_lefts[j + 1]:
            c = j
            break
    return (r, c) if r is not None and c is not None else None


_MAX_XLSX_ROWS = 500


def _excel_ole_segment(shape, *, slide_no, warnings):
    """엑셀에서 붙여넣은 임베드(OLE) 표 → table 세그먼트.

    PPTX 안에는 원본 파일이 통째로 들어있어(`ole_format.blob`) 신형 .xlsx 는
    openpyxl 로 열어 표로 읽는다. 구형 .xls(OLE 복합문서)는 못 읽어 안내만 남긴다.
    """
    of = getattr(shape, "ole_format", None)
    if of is None:
        return None
    try:
        prog = of.prog_id or ""
        blob = of.blob
    except Exception:
        return None
    if not blob:
        return None
    if not (prog.startswith("Excel") or blob[:2] == b"PK"):
        # 엑셀이 아닌 임베드 개체(워드 문서 등) — 못 읽음
        warnings.append(
            f"슬라이드 {slide_no}: 지원 안 되는 임베드 개체({prog or '알 수 없음'}) 건너뜀"
        )
        return None
    if blob[:2] != b"PK":
        warnings.append(
            f"슬라이드 {slide_no}: 구형 엑셀(.xls) 임베드 표는 못 읽음 "
            "— PPT 「표 삽입」으로 다시 넣어주세요"
        )
        return None
    try:
        import openpyxl

        wb = openpyxl.load_workbook(BytesIO(blob), read_only=True, data_only=True)
        ws = wb.active
        rows = []
        truncated = False
        for row in ws.iter_rows(values_only=True):
            rows.append(["" if v is None else str(v).strip() for v in row])
            if len(rows) >= _MAX_XLSX_ROWS:
                truncated = True
                break
        wb.close()
    except Exception:
        warnings.append(f"슬라이드 {slide_no}: 엑셀 임베드 표 읽기 실패 1건")
        return None
    # 빈 행 제거 + 우측 빈 열 트림
    rows = [r for r in rows if any(c for c in r)]
    if not rows:
        return None
    width = 0
    for r in rows:
        last = 0
        for i, c in enumerate(r):
            if c:
                last = i + 1
        width = max(width, last)
    rows = [r[:width] for r in rows]
    if truncated:
        warnings.append(
            f"슬라이드 {slide_no}: 엑셀 임베드 표가 커서 {_MAX_XLSX_ROWS}행까지만 가져옴"
        )
    return _table_segment(rows)


def _shape_to_segment(shape, *, is_title, slide_no, warnings):
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    stype = None
    try:
        stype = shape.shape_type
    except Exception:
        stype = None

    # 그림
    if stype == MSO_SHAPE_TYPE.PICTURE:
        try:
            img = shape.image
            return {
                "type": "image",
                "image_bytes": img.blob,
                "image_mime": img.content_type or "image/png",
                "image_ext": (img.ext or "png").lstrip("."),
                "alt": (getattr(shape, "name", "") or ""),
            }
        except Exception:
            warnings.append(f"슬라이드 {slide_no}: 이미지 추출 실패 1건")
            return None

    # 표
    if getattr(shape, "has_table", False):
        try:
            return _table_segment(_table_grid(shape.table, overlay=None))
        except Exception:
            warnings.append(f"슬라이드 {slide_no}: 표 추출 실패 1건")
            return None

    # 차트(미지원) — 명시적으로 건너뜀 기록
    if getattr(shape, "has_chart", False):
        warnings.append(f"슬라이드 {slide_no}: 차트 1개 건너뜀(이미지로 다시 붙여주세요)")
        return None

    # 텍스트
    if getattr(shape, "has_text_frame", False):
        try:
            tf = shape.text_frame
            if is_title:
                text = (tf.text or "").strip()
                return {"type": "heading", "content": {"text": text, "level": 1}} if text else None
            items = []
            for p in tf.paragraphs:
                t = (p.text or "").strip()
                if not t:
                    continue
                depth = p.level if isinstance(p.level, int) else 0
                items.append({"depth": max(0, min(MAX_DEPTH, depth)), "text": t})
            return {"type": "rich_text", "content": {"items": items}} if items else None
        except Exception:
            return None

    # 엑셀에서 붙여넣은 임베드(OLE) 표 → 표로 읽어들임
    seg = _excel_ole_segment(shape, slide_no=slide_no, warnings=warnings)
    if seg:
        return seg

    # 어디에도 안 걸린 개체 — 유형을 warnings 에 남긴다(SmartArt·구형 임베드 등이
    # "표를 표로 못 읽는" 흔한 원인). 장식/빈 도형은 흔하니 제외해 소음을 줄인다.
    tname = getattr(stype, "name", "") or ""
    # EMBEDDED_OLE_OBJECT 는 위 _excel_ole_segment 가 전담(자체 안내) → 중복 제외.
    _NOISE = {
        "AUTO_SHAPE", "TEXT_BOX", "PLACEHOLDER", "LINE",
        "FREEFORM", "CONNECTOR", "EMBEDDED_OLE_OBJECT",
    }
    if tname and tname not in _NOISE:
        label = tname
        if tname == "TABLE":  # 방어적: 여기 오면 안 되지만
            label = "표(추출 실패)"
        warnings.append(f"슬라이드 {slide_no}: 지원 안 되는 개체({label}) 건너뜀")
    return None


def parse_pptx(data: bytes) -> dict:
    from pptx import Presentation

    prs = Presentation(BytesIO(data))
    pages = []
    warnings: list[str] = []
    for si, slide in enumerate(prs.slides, 1):
        try:
            title_shape = slide.shapes.title
        except Exception:
            title_shape = None
        name = ""
        # python-pptx 는 접근할 때마다 새 래퍼 객체를 주므로 `is` 비교가 안 된다 →
        # 안정적인 shape_id 로 제목을 식별한다.
        title_id = None
        if title_shape is not None:
            try:
                name = (title_shape.text_frame.text or "").strip()
            except Exception:
                name = ""
            try:
                title_id = title_shape.shape_id
            except Exception:
                title_id = None
        name = name[:120] or f"슬라이드 {si}"

        shapes = _flatten(slide.shapes)
        shapes.sort(key=_pos)

        # 표 위에 겹쳐 놓은 값 텍스트박스를 셀로 흡수한다(표지 양식이 흔한 패턴:
        # 표는 라벨+빈 값칸, 실제 값은 위에 얹은 텍스트박스). 각 텍스트박스의
        # 좌상단이 어느 표의 빈 셀에 들어가면 그 셀 값으로 넣고 별도 문단으로는
        # 내보내지 않는다.
        tables = []  # (shape_id, shape, row_tops, col_lefts, bbox, overlay)
        for shape in shapes:
            if not getattr(shape, "has_table", False):
                continue
            try:
                row_tops, col_lefts = _table_geom(shape)
                bbox = (row_tops[0], col_lefts[0], row_tops[-1], col_lefts[-1])
                tables.append([shape.shape_id, shape, row_tops, col_lefts, bbox, {}])
            except Exception:
                pass

        consumed = set()
        if tables:
            for shape in shapes:
                try:
                    sid = shape.shape_id
                except Exception:
                    continue
                if sid == title_id or getattr(shape, "has_table", False):
                    continue
                if not getattr(shape, "has_text_frame", False):
                    continue
                txt = (shape.text_frame.text or "").strip()
                if not txt:
                    continue
                top, left = _pos(shape)
                for tid, tshape, row_tops, col_lefts, bbox, overlay in tables:
                    t0, l0, t1, l1 = bbox
                    if not (t0 <= top < t1 and l0 <= left < l1):
                        continue
                    rc = _locate_cell(row_tops, col_lefts, top, left)
                    if rc is None:
                        continue
                    r, c = rc
                    # 이미 라벨이 있는 셀은 건드리지 않고 문단으로 남긴다.
                    if (tshape.table.cell(r, c).text or "").strip():
                        continue
                    overlay.setdefault((r, c), txt)
                    consumed.add(sid)
                    break

        overlay_by_table = {t[0]: t[5] for t in tables}

        segments = []
        for shape in shapes:
            try:
                sid = shape.shape_id
            except Exception:
                sid = None
            if sid is not None and sid in consumed:
                continue
            if getattr(shape, "has_table", False):
                try:
                    seg = _table_segment(
                        _table_grid(shape.table, overlay=overlay_by_table.get(sid))
                    )
                except Exception:
                    warnings.append(f"슬라이드 {si}: 표 추출 실패 1건")
                    seg = None
            else:
                is_title = title_id is not None and sid == title_id
                seg = _shape_to_segment(
                    shape, is_title=is_title, slide_no=si, warnings=warnings
                )
            if seg:
                segments.append(seg)
        pages.append({"name": name, "segments": segments})
    return {"pages": pages, "warnings": warnings}
